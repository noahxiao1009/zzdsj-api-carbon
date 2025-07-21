"""
文本提取服务
从不同格式的文档中提取文本内容
"""

import os
import asyncio
from pathlib import Path
from typing import Dict, Any, Optional
import logging

# 文档处理库
try:
    import PyPDF2
    import pdfplumber
except ImportError:
    PyPDF2 = None
    pdfplumber = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

logger = logging.getLogger(__name__)


class TextExtractor:
    """文本提取服务"""
    
    def __init__(self):
        self.supported_formats = {
            'pdf': self._extract_pdf,
            'word': self._extract_word,
            'text': self._extract_text,
            'markdown': self._extract_text,
            'csv': self._extract_csv,
            'excel': self._extract_excel,
            'powerpoint': self._extract_powerpoint,
            'html': self._extract_html,
            'xml': self._extract_xml,
            'json': self._extract_json
        }
    
    async def extract_text(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """
        从文档中提取文本
        
        Args:
            file_path: 文件路径
            file_type: 文件类型
            
        Returns:
            包含提取结果的字典
        """
        try:
            if file_type not in self.supported_formats:
                return {
                    'success': False,
                    'error': f'Unsupported file type: {file_type}'
                }
            
            # 检查文件是否存在
            if not Path(file_path).exists():
                return {
                    'success': False,
                    'error': f'File not found: {file_path}'
                }
            
            # 调用对应的提取方法
            extractor = self.supported_formats[file_type]
            result = await extractor(file_path)
            
            if result['success']:
                # 生成内容预览
                content = result.get('content', '')
                result['content_preview'] = self._generate_preview(content)
                result['char_count'] = len(content)
                result['word_count'] = len(content.split()) if content else 0
            
            return result
            
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return {
                'success': False,
                'error': f'Text extraction failed: {str(e)}'
            }
    
    async def _extract_pdf(self, file_path: str) -> Dict[str, Any]:
        """提取PDF文档文本"""
        if not pdfplumber and not PyPDF2:
            return {
                'success': False,
                'error': 'PDF processing libraries not available'
            }
        
        try:
            content = ""
            metadata = {}
            
            # 优先使用pdfplumber，提取效果更好
            if pdfplumber:
                with pdfplumber.open(file_path) as pdf:
                    metadata = {
                        'page_count': len(pdf.pages),
                        'pdf_metadata': pdf.metadata or {}
                    }
                    
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            content += page_text + "\n"
            
            # 备用PyPDF2
            elif PyPDF2:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    metadata = {
                        'page_count': len(pdf_reader.pages),
                        'pdf_metadata': pdf_reader.metadata or {}
                    }
                    
                    for page in pdf_reader.pages:
                        content += page.extract_text() + "\n"
            
            return {
                'success': True,
                'content': content.strip(),
                'metadata': metadata,
                'extraction_method': 'pdfplumber' if pdfplumber else 'PyPDF2'
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'PDF extraction failed: {str(e)}'
            }
    
    async def _extract_word(self, file_path: str) -> Dict[str, Any]:
        """提取Word文档文本"""
        if not DocxDocument:
            return {
                'success': False,
                'error': 'Word processing library not available'
            }
        
        try:
            doc = DocxDocument(file_path)
            
            # 提取段落文本
            paragraphs = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text)
            
            content = "\n".join(paragraphs)
            
            # 提取表格文本
            tables_text = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(" | ".join(row_data))
                tables_text.append("\n".join(table_data))
            
            if tables_text:
                content += "\n\n" + "\n\n".join(tables_text)
            
            metadata = {
                'paragraph_count': len(paragraphs),
                'table_count': len(doc.tables)
            }
            
            return {
                'success': True,
                'content': content,
                'metadata': metadata
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'Word extraction failed: {str(e)}'
            }
    
    async def _extract_text(self, file_path: str) -> Dict[str, Any]:
        """提取纯文本文件"""
        try:
            # 尝试不同编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                    
                    return {
                        'success': True,
                        'content': content,
                        'metadata': {'encoding': encoding}
                    }
                except UnicodeDecodeError:
                    continue
            
            return {
                'success': False,
                'error': 'Unable to decode text file with supported encodings'
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'Text extraction failed: {str(e)}'
            }
    
    async def _extract_csv(self, file_path: str) -> Dict[str, Any]:
        """提取CSV文件内容"""
        if not pd:
            return {
                'success': False,
                'error': 'Pandas library not available'
            }
        
        try:
            # 尝试不同编码读取CSV
            encodings = ['utf-8', 'gbk', 'gb2312']
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                return {
                    'success': False,
                    'error': 'Unable to read CSV with supported encodings'
                }
            
            # 将DataFrame转换为文本
            content_lines = []
            
            # 添加列标题
            content_lines.append(" | ".join(df.columns.astype(str)))
            
            # 添加数据行
            for _, row in df.iterrows():
                content_lines.append(" | ".join(row.astype(str)))
            
            content = "\n".join(content_lines)
            
            metadata = {
                'row_count': len(df),
                'column_count': len(df.columns),
                'columns': df.columns.tolist()
            }
            
            return {
                'success': True,
                'content': content,
                'metadata': metadata
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'CSV extraction failed: {str(e)}'
            }
    
    async def _extract_excel(self, file_path: str) -> Dict[str, Any]:
        """提取Excel文件内容"""
        if not pd:
            return {
                'success': False,
                'error': 'Pandas library not available'
            }
        
        try:
            # 读取所有工作表
            excel_file = pd.ExcelFile(file_path)
            all_content = []
            metadata = {'sheet_names': excel_file.sheet_names, 'sheets': {}}
            
            for sheet_name in excel_file.sheet_names:
                df = excel_file.parse(sheet_name)
                
                # 添加工作表标题
                all_content.append(f"=== {sheet_name} ===")
                
                # 添加列标题
                all_content.append(" | ".join(df.columns.astype(str)))
                
                # 添加数据行
                for _, row in df.iterrows():
                    all_content.append(" | ".join(row.astype(str)))
                
                all_content.append("")  # 空行分隔
                
                metadata['sheets'][sheet_name] = {
                    'row_count': len(df),
                    'column_count': len(df.columns)
                }
            
            content = "\n".join(all_content)
            
            return {
                'success': True,
                'content': content,
                'metadata': metadata
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'Excel extraction failed: {str(e)}'
            }
    
    async def _extract_powerpoint(self, file_path: str) -> Dict[str, Any]:
        """提取PowerPoint文件内容"""
        if not Presentation:
            return {
                'success': False,
                'error': 'PowerPoint processing library not available'
            }
        
        try:
            prs = Presentation(file_path)
            content_parts = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_content = [f"=== Slide {i} ==="]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                content_parts.append("\n".join(slide_content))
            
            content = "\n\n".join(content_parts)
            
            metadata = {
                'slide_count': len(prs.slides)
            }
            
            return {
                'success': True,
                'content': content,
                'metadata': metadata
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'PowerPoint extraction failed: {str(e)}'
            }
    
    async def _extract_html(self, file_path: str) -> Dict[str, Any]:
        """提取HTML文件内容"""
        if not BeautifulSoup:
            return {
                'success': False,
                'error': 'BeautifulSoup library not available'
            }
        
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 移除script和style标签
            for script in soup(["script", "style"]):
                script.decompose()
            
            # 提取文本
            text = soup.get_text()
            
            # 清理文本
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            content = '\n'.join(chunk for chunk in chunks if chunk)
            
            metadata = {
                'title': soup.title.string if soup.title else None,
                'has_tables': len(soup.find_all('table')) > 0
            }
            
            return {
                'success': True,
                'content': content,
                'metadata': metadata
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'HTML extraction failed: {str(e)}'
            }
    
    async def _extract_xml(self, file_path: str) -> Dict[str, Any]:
        """提取XML文件内容"""
        if not BeautifulSoup:
            return {
                'success': False,
                'error': 'BeautifulSoup library not available'
            }
        
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                xml_content = file.read()
            
            soup = BeautifulSoup(xml_content, 'xml')
            content = soup.get_text()
            
            return {
                'success': True,
                'content': content,
                'metadata': {'format': 'xml'}
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'XML extraction failed: {str(e)}'
            }
    
    async def _extract_json(self, file_path: str) -> Dict[str, Any]:
        """提取JSON文件内容"""
        try:
            import json
            
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
            
            # 将JSON转换为可读文本
            content = json.dumps(data, ensure_ascii=False, indent=2)
            
            return {
                'success': True,
                'content': content,
                'metadata': {'format': 'json'}
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'JSON extraction failed: {str(e)}'
            }
    
    def _generate_preview(self, content: str, max_length: int = 500) -> str:
        """生成内容预览"""
        if not content:
            return ""
        
        if len(content) <= max_length:
            return content
        
        # 截取前max_length个字符，并在单词边界处截断
        preview = content[:max_length]
        last_space = preview.rfind(' ')
        
        if last_space > max_length * 0.8:  # 如果最后一个空格位置合理
            preview = preview[:last_space]
        
        return preview + "..."